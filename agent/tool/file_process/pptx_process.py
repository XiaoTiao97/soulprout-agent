from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from datetime import datetime

def extract_ppt_complete_info(ppt_path):
    """
    提取PPT文件的完整结构化信息（修复版）
    """
    try:
        prs = Presentation(ppt_path)

        ppt_info = {
            'file_info': {
                'slide_count': len(prs.slides),
                'slide_size': {
                    'width': prs.slide_width,
                    'height': prs.slide_height
                },
                'extraction_time': datetime.now().isoformat()
            },
            'masters': [],
            'layouts': [],
            'slides': []
        }

        # 提取母版和布局信息
        for master in prs.slide_masters:
            master_info = {
                'master_id': len(ppt_info['masters']),
                'shape_count': len(master.shapes)
            }
            ppt_info['masters'].append(master_info)

        for layout in prs.slide_layouts:
            layout_info = {
                'layout_id': len(ppt_info['layouts']),
                'name': layout.name,
                'placeholder_count': sum(1 for shape in layout.shapes if shape.is_placeholder)
            }
            ppt_info['layouts'].append(layout_info)

        # 提取幻灯片信息
        for slide_index, slide in enumerate(prs.slides):
            slide_info = {
                'slide_number': slide_index + 1,
                'layout': slide.slide_layout.name,
                'shapes': []
            }

            for shape in slide.shapes:
                shape_info = extract_shape_info_fixed(shape)
                if shape_info:
                    slide_info['shapes'].append(shape_info)

            ppt_info['slides'].append(slide_info)

        return str(ppt_info)

    except Exception as e:
        return {'error': f'提取失败: {str(e)}'}


def extract_shape_info_fixed(shape):
    """修复版形状信息提取"""
    try:
        shape_info = {
            'shape_id': getattr(shape, 'shape_id', None),
            'name': getattr(shape, 'name', ''),
            'type': get_shape_type_name(shape),
            'position': {
                'left': int(getattr(shape, 'left', 0)),
                'top': int(getattr(shape, 'top', 0)),
                'width': int(getattr(shape, 'width', 0)),
                'height': int(getattr(shape, 'height', 0))
            }
        }

        # 关键修复：提取文本内容（包括标题）
        if hasattr(shape, 'text') and shape.text.strip():
            shape_info['text_content'] = shape.text

        # 根据形状类型提取特定信息
        shape_type = getattr(shape, 'shape_type', None)

        if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            shape_info.update(extract_text_info_fixed(shape))
        elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            shape_info.update(extract_placeholder_info(shape))
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_info.update(extract_picture_info_fixed(shape))
        elif shape_type == MSO_SHAPE_TYPE.TABLE:
            shape_info.update(extract_table_info_fixed(shape))
        elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_info.update(extract_autoshape_info_fixed(shape))
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            shape_info.update(extract_group_info_fixed(shape))

        return shape_info

    except Exception as e:
        return {
            'shape_id': getattr(shape, 'shape_id', None),
            'name': getattr(shape, 'name', ''),
            'type': 'UNKNOWN',
            'error': f'提取失败: {str(e)}'
        }


def extract_text_info_fixed(shape):
    """修复版文本信息提取 - 合并相同格式的文本段"""
    text_info = {
        'content_type': 'text',
        'text_content': getattr(shape, 'text', '')
    }

    try:
        if hasattr(shape, 'text_frame'):
            paragraphs = []

            for para in shape.text_frame.paragraphs:
                para_text = getattr(para, 'text', '')
                if para_text.strip():
                    para_info = {
                        'text': para_text,
                        'alignment': get_alignment_name(para.alignment) if hasattr(para, 'alignment') else None,
                        'font_runs': merge_font_runs(para.runs)  # 关键修改：合并相同格式的文本段
                    }

                    paragraphs.append(para_info)

            if paragraphs:
                text_info['paragraphs'] = paragraphs

    except Exception as e:
        text_info['text_error'] = f'文本格式提取失败: {str(e)}'

    return text_info


def merge_font_runs(runs):
    """合并相同字体格式的连续文本段"""
    if not runs:
        return []

    merged_runs = []
    current_run = {
        'text': runs[0].text,
        'font': extract_font_info(runs[0].font) if hasattr(runs[0], 'font') else {}
    }

    for i in range(1, len(runs)):
        current_font = extract_font_info(runs[i].font) if hasattr(runs[i], 'font') else {}

        # 如果字体格式相同，合并文本
        if current_run['font'] == current_font:
            current_run['text'] += runs[i].text
        else:
            # 格式不同，保存当前段，开始新段
            merged_runs.append(current_run)
            current_run = {
                'text': runs[i].text,
                'font': current_font
            }

    # 添加最后一个文本段
    merged_runs.append(current_run)

    return merged_runs


def extract_placeholder_info(shape):
    """占位符信息提取"""
    placeholder_info = {
        'content_type': 'placeholder'
    }

    try:
        if hasattr(shape, 'placeholder_format'):
            placeholder_format = shape.placeholder_format
            placeholder_info['placeholder_type'] = get_placeholder_type_name(
                getattr(placeholder_format, 'type', None)
            )

        # 占位符也有文本内容
        if hasattr(shape, 'text') and shape.text.strip():
            placeholder_info['text_content'] = shape.text

            # 提取占位符的字体信息
            if hasattr(shape, 'text_frame'):
                placeholder_info.update(extract_text_info_fixed(shape))

    except Exception as e:
        placeholder_info['placeholder_error'] = str(e)

    return placeholder_info


def extract_font_info(font):
    """提取字体信息（安全版）"""
    font_info = {}

    try:
        # 基本字体属性
        if hasattr(font, 'name') and font.name:
            font_info['name'] = font.name

        if hasattr(font, 'size') and font.size:
            font_info['size'] = int(font.size) if font.size else None

        if hasattr(font, 'bold'):
            font_info['bold'] = font.bold

        if hasattr(font, 'italic'):
            font_info['italic'] = font.italic

        if hasattr(font, 'underline'):
            font_info['underline'] = font.underline

        # 颜色信息（安全提取）
        if hasattr(font, 'color'):
            color_info = extract_color_info(font.color)
            if color_info:
                font_info['color'] = color_info

    except Exception:
        pass  # 忽略单个字体属性提取失败

    return font_info


def extract_color_info(color_obj):
    """安全提取颜色信息"""
    if color_obj is None:
        return None

    try:
        if hasattr(color_obj, 'rgb') and color_obj.rgb:
            return f"#{color_obj.rgb:06X}"
        elif hasattr(color_obj, 'theme_color'):
            return f"theme_{color_obj.theme_color}"
        elif hasattr(color_obj, 'type'):
            return f"color_type_{color_obj.type}"
    except Exception:
        pass

    return None


def extract_picture_info_fixed(shape):
    """修复版图片信息提取"""
    return {
        'content_type': 'picture',
        'cropped': hasattr(shape, 'crop_left') and any([
            getattr(shape, 'crop_left', 0) > 0,
            getattr(shape, 'crop_right', 0) > 0,
            getattr(shape, 'crop_top', 0) > 0,
            getattr(shape, 'crop_bottom', 0) > 0
        ])
    }


def extract_autoshape_info_fixed(shape):
    """修复版自动形状信息提取"""
    shape_info = {
        'content_type': 'shape'
    }

    try:
        # 形状的文本内容
        if hasattr(shape, 'text') and shape.text.strip():
            shape_info['text_content'] = shape.text
            shape_info.update(extract_text_info_fixed(shape))
    except Exception:
        pass

    return shape_info


def extract_table_info_fixed(shape):
    """修复版表格信息提取"""
    try:
        table = shape.table
        table_info = {
            'content_type': 'table',
            'dimensions': {
                'rows': len(table.rows),
                'columns': len(table.columns)
            },
            'cells': []
        }

        # 提取表格内容
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_text = getattr(cell, 'text', '').strip()
                if cell_text:
                    table_info['cells'].append({
                        'row': row_idx,
                        'column': col_idx,
                        'text': cell_text
                    })

        return table_info
    except Exception:
        return {'content_type': 'table', 'error': '表格信息提取失败'}


def extract_group_info_fixed(shape):
    """修复版组合形状信息提取"""
    group_info = {
        'content_type': 'group',
        'child_count': len(shape.shapes) if hasattr(shape, 'shapes') else 0
    }

    try:
        # 组合形状的文本内容
        if hasattr(shape, 'text') and shape.text.strip():
            group_info['text_content'] = shape.text
    except Exception:
        pass

    return group_info


def get_placeholder_type_name(placeholder_type):
    """获取占位符类型的友好名称"""
    if placeholder_type is None:
        return None

    type_map = {
        1: 'TITLE',
        2: 'BODY',
        3: 'CENTER_TITLE',
        4: 'SUBTITLE',
        7: 'OBJECT',
        13: 'SLIDE_NUMBER',
        14: 'PLACEHOLDER',
        15: 'FOOTER',
        16: 'DATE',
        18: 'PICTURE'
    }

    return type_map.get(placeholder_type, f'PLACEHOLDER_{placeholder_type}')


# 其他辅助函数保持不变...
def get_shape_type_name(shape):
    """获取形状类型的友好名称"""
    try:
        shape_type = getattr(shape, 'shape_type', None)
        if shape_type is None:
            return 'UNKNOWN'

        type_map = {
            MSO_SHAPE_TYPE.TEXT_BOX: 'TEXT_BOX',
            MSO_SHAPE_TYPE.PICTURE: 'PICTURE',
            MSO_SHAPE_TYPE.TABLE: 'TABLE',
            MSO_SHAPE_TYPE.AUTO_SHAPE: 'SHAPE',
            MSO_SHAPE_TYPE.GROUP: 'GROUP',
            MSO_SHAPE_TYPE.PLACEHOLDER: 'PLACEHOLDER',
            MSO_SHAPE_TYPE.LINE: 'LINE'
        }

        return type_map.get(shape_type, str(shape_type))
    except Exception:
        return 'UNKNOWN'


def get_alignment_name(alignment):
    """获取对齐方式的友好名称"""
    if alignment is None:
        return None

    alignment_map = {
        0: 'LEFT',
        1: 'CENTER',
        2: 'RIGHT',
        3: 'JUSTIFY'
    }

    return alignment_map.get(alignment, str(alignment))


# 使用示例
if __name__ == "__main__":
    # 使用方法
    result = extract_ppt_complete_info("Agent能力.pptx")
    from pprint import pprint
    pprint(result)

